import asyncio
import logging
from concurrent.futures import ThreadPoolExecutor
from io import BytesIO

from pptx import Presentation
from pptx.util import Pt

from app.extractors.base import BaseExtractor

logger = logging.getLogger(__name__)

SKIP_PLACEHOLDER_TYPES = {13, 14}  # slide number, date placeholders


class PptxExtractor(BaseExtractor):
    """
    Extracts text from .pptx files using python-pptx.
    Covers: text frames, tables, and speaker notes per slide.
    """

    def __init__(self, executor: ThreadPoolExecutor):
        self._executor = executor

    async def extract(self, content: bytes, filename: str = "") -> str:
        loop = asyncio.get_event_loop()
        return await loop.run_in_executor(
            self._executor, self._sync_extract, content
        )

    def _sync_extract(self, content: bytes) -> str:
        prs = Presentation(BytesIO(content))
        slides_text = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            parts = [f"## Slide {slide_num}"]

            # Extract text from all shapes on the slide
            for shape in slide.shapes:
                # Skip slide number / date placeholders
                if (
                    shape.is_placeholder
                    and shape.placeholder_format.idx in SKIP_PLACEHOLDER_TYPES
                ):
                    continue

                # Text frame shapes (titles, content boxes, text boxes)
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            parts.append(text)

                # Table shapes
                elif shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        cells = [
                            cell.text.strip() for cell in row.cells
                        ]
                        # Skip rows where all cells are empty
                        if any(cells):
                            parts.append(" | ".join(cells))

            # Speaker notes (often contain research context and citations)
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    parts.append(f"\n_Notes: {notes_text}_")

            slides_text.append("\n".join(parts))

        return "\n\n---\n\n".join(slides_text)
